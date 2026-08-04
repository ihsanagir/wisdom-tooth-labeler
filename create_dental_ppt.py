"""
coTalent Cyberpark Sunumu — Dental Proje Versiyonu
Kaynak PPT'yi template olarak kullanır, tüm başlıklar/logolar korunur,
sadece içerik dental projeyle doldurulur.
"""

import copy
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

SRC = r'C:\Users\İhsan Furkan\Desktop\coTalent_FirmaSunumu_14012026_yedek.pptx'
DST = r'C:\Users\İhsan Furkan\Desktop\Dental_Yirmilik_Dis_Cyberpark_Sunumu.pptx'

prs = Presentation(SRC)
slides = list(prs.slides)


def safe_set(slide, name, text):
    """Placeholder'ı name ile bul ve metni temiz sekilde degistir."""
    # Kontrol karakterlerini temizle (NULL, VT vb)
    text = text.replace('\x00', '').replace('\x0b', '\n').replace('\r', '')
    
    for shape in slide.shapes:
        if shape.name == name and hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            txBody = tf._txBody
            paras = txBody.findall(qn('a:p'))

            # Ilk paragrafin formatlama ozellikleri
            saved_pPr = None
            saved_rPr = None
            if paras:
                saved_pPr = paras[0].find(qn('a:pPr'))
                runs = paras[0].findall(qn('a:r'))
                if runs:
                    saved_rPr = runs[0].find(qn('a:rPr'))

            # Tum a:p'leri kaldir
            for p in paras:
                txBody.remove(p)

            # Satirlara bol ve yeni a:p ekle
            lines = text.split('\n')
            for li, line in enumerate(lines):
                p_new = etree.SubElement(txBody, qn('a:p'))
                if li == 0 and saved_pPr is not None:
                    p_new.insert(0, copy.deepcopy(saved_pPr))
                r_new = etree.SubElement(p_new, qn('a:r'))
                if saved_rPr is not None:
                    r_new.insert(0, copy.deepcopy(saved_rPr))
                t_new = etree.SubElement(r_new, qn('a:t'))
                t_new.text = line
            return True
    return False


# ===================================================================
# SLIDE 1 — Kapak
# ===================================================================
s = slides[0]
safe_set(s, 'PlaceHolder 1',
    'DEGERLENDIRME KOMISYONU SUNUMU\n'
    'Derin Ogrenme Tabanli Goruntu Isleme Yontemleriyle\n'
    'Gomulu Yirmilik Dis Tespiti ve Cerrahi Risk Siniflandirmasi\n'
    '\n'
    '\n'
    '19.06.2025\n'
    'Revizyon Tarihi: 19.06.2025\n'
    'Hazirlayan : Ihsan Furkan AGIR  Bilgisayar Muhendisi'
)

# ===================================================================
# SLIDE 2 — Sirket Genel Bilgileri
# ===================================================================
s = slides[1]
safe_set(s, 'PlaceHolder 2',
    'Aktarilan bilgiler Basvuru Raporu\'nda aktarilan bilgi ile uyumlu olmalidir.\n'
    'Unvani: coTalent Bilisim Teknolojileri Sanayi Ticaret Ltd. Sti.\n'
    'Personel Sayisi: (1)\n'
    '\n'
    '           Ihsan Furkan AGIR - Founder (Bilgisayar Muhendisi)\n'
    '           Full Stack Development, Makine Ogrenmesi, Goruntu Isleme, AI Modelling\n'
    '\n'
    'Danishman Sayisi: (3)\n'
    '\n'
    'Goruntu Isleme ve Bilgisayarli Goru Danismani,\n'
    ' Medikal Goruntu Analizi, Nesne Tespiti, Derin Ogrenme\n'
    '\n'
    'Dis Hekimligi Klinik Danismani\n'
    'Agiz, Dis ve Cene Cerrahisi, Pell-Gregory Siniflandirmasi, Klinik Protokoller\n'
    '\n'
    'Yapay Zeka Danismani\n'
    'Bilgisayar Gorusu, Model Optimizasyonu, ML/Ops'
)

# ===================================================================
# SLIDE 3 — Sirket Ortaklik Yapisi
# ===================================================================
s = slides[2]
safe_set(s, 'PlaceHolder 2',
    'Ortak sayisi: (1) Ihsan Furkan AGIR\n'
    'Her ortak icin ayri ayri doldurulacak bilgiler:\n'
    'Sirketteki gorevi: Kurucu / Yazilim Gelistirici\n'
    'Sirketteki payi: 100%\n'
    'Kisa ozgecmisi: Bilgisayar Muhendisi\n'
    '\n'
    'Firmanin standartlari: Gelistirme surecinde ISO 13485 Medikal Yazilim\n'
    'standartlarina uyum hedeflenmektedir.\n'
    '\n'
    'Firmanin referanslari:\n'
    '1. Yapay Zeka destekli dis projesi PoC (Proof of Concept)\n'
    '2. Web tabanli panoramik rontgen analiz yazilimi prototipi\n'
    '3. YOLO tabanli medikal goruntu tespit modeli (disprojesi3)\n'
)

# ===================================================================
# SLIDE 4 — Sirket Ciro Bilgisi (aynen birak - yeni kurulan sirket)
# ===================================================================
# Degisiklik yok.

# ===================================================================
# SLIDE 5 — Proje Bilgileri: Giris / Problem Tanimi
# ===================================================================
s = slides[4]
safe_set(s, 'PlaceHolder 3',
    '"Agiz, dis ve cene cerrahisinde yirmilik dis (ucuncu molar) operasyonlari\n'
    'en sik gerceklestirilen cerrahi mudahaleler arasindadir."\n'
    '\n'
    '"Panoramik rontgen degerlendirmesi deneyime bagli,\n'
    'zaman alici ve gozlemciden gozlemciye tutarsizdir."\n'
    '\n'
    '"Mevcut degerlendirme yontemleri\n'
    'manuel / subjektif / tecrube bagimlI / hata payi yuksek."\n'
)

# ===================================================================
# SLIDE 6 — Problem Istatistikleri
# ===================================================================
s = slides[5]
safe_set(s, 'PlaceHolder 1', 'Proje Bilgileri : Problem Istatistikleri')
safe_set(s, 'PlaceHolder 3',
    'Turkiye\'de yillik yaklasik 2-3 milyon yirmilik dis operasyonu gerceklestirilmektedir.\n'
    '\n'
    'Pell-Gregory ve Winter siniflandirma sistemleri hala subjektif degerlendirmeye\n'
    'dayanmakta; farkli klinisyenler ayni rontgen icin farkli siniflar atayabilmektedir.\n'
    '\n'
    'Yanlis zorluk tahmini; beklenmedik komplikasyonlara (sinir hasari,\n'
    'alveolar kirik, kok kirigi), uzun operasyon surelerine ve\n'
    'hasta memnuniyetsizligine yol acmaktadir.\n'
    '\n'
    'Yapay zeka destekli goruntu analizi, tani dogruluğunu artirma\n'
    've degerlendirme suresini dramatik bicimde kisaltma potansiyeline sahiptir.\n'
)

# ===================================================================
# SLIDE 7 — Problem: Komplikasyon Istatistikleri
# ===================================================================
s = slides[6]
safe_set(s, 'PlaceHolder 1', 'Problem Tanimi: Komplikasyon Istatistikleri')
safe_set(s, 'PlaceHolder 3',
    'Inferior alveoler sinir (IAN) hasari, en ciddi yirmilik dis cerrahisi komplikasyonudur.\n'
    'Risk; gomululuk tipi, derinlik ve sinir komsulugˆuyla dogrudan iliskilidir.\n'
    '\n'
    'Istatistikler gostermektedir ki cerrahi oncesi yetersiz radyolojik degerlendirme,\n'
    'komplikasyon oranini 3-5 kat artirmaktadir.\n'
    '\n'
    'Yatay (Horizontal) gomulu dislerde sinir hasari riski dikey (Vertical)\n'
    'gomulu dislere kiyasla yaklasik 4 kat daha yuksektir.\n'
    '\n'
    'Dogru on degerlendirme => dogru tedavi plani => komplikasyon riski azalir.\n'
)

# ===================================================================
# SLIDE 8 — Problem: Klinik Risk Verileri (gorsel slayt)
# ===================================================================
s = slides[7]
safe_set(s, 'PlaceHolder 1', 'Problem Tanimi: Klinik Risk Verileri')

# ===================================================================
# SLIDE 9 — Problem: Mevcut Degerlendirme Yontemleri
# ===================================================================
s = slides[8]
safe_set(s, 'PlaceHolder 1', 'Problem Tanimi: Mevcut Degerlendirme Yontemleri (Geleneksel Gozlem)')
safe_set(s, 'PlaceHolder 3',
    '"Manuel panoramik rontgen degerlendirmesi\n'
    '=> zaman alici, subjektif, gozlemciye bagimlI."\n'
    '\n'
    '"Pell-Gregory / Winter siniflandirmasi\n'
    '=> tutarsiz, tecrube gerektiren, standart disi."\n'
    '\n'
    '"CBCT ile 3D degerlendirme\n'
    '=> pahali, yuksek radyasyon dozu, her klinikte mevcut degil."\n'
    '\n'
    'Yanlis cerrahi zorluk tahmini genellikle fark edilmeden operasyon\n'
    'esnasinda ortaya cikar.\n'
    '\n'
    'Vakalarin onemli bir kisminda detayli on degerlendirme yapilmamaktadir.\n'
)

# ===================================================================
# SLIDE 10 — Teklif Edilen Cozum
# ===================================================================
s = slides[9]
safe_set(s, 'PlaceHolder 3',
    'Panoramik rontgen goruntusu uzerinden otomatik dis tespiti & parametre cikarimi\n'
    '\n'
    'Derin Ogrenme (YOLO) ile gomulu yirmilik dis konumlandirmasi\n'
    '\n'
    '"Yapay zeka + goruntu isleme + kural tabanli skorlama motoru ile\n'
    'panoramik rontgenden saniyeler icinde cerrahi zorluk tahmini."\n'
    '\n'
    'Modern yaklasimlar (YOLOv8, YOLO11)\n'
    'ile yuksek dogruluk, standardizasyon ve olceklenebilirlik mumkun.\n'
)

# ===================================================================
# SLIDE 11 — Teklif Edilen Cozum: Hedeflenen Teknolojik Kazanimlar
# ===================================================================
s = slides[10]
safe_set(s, 'PlaceHolder 1', 'Teklif Edilen Cozum : Hedeflenen Teknolojik Kazanimlar')
safe_set(s, 'PlaceHolder 3',
    '- YOLOv8/YOLO11 tabanli nesne tespiti modeli\n'
    '  => Panoramik rontgenden gomulu disleri otomatik konumlandirma\n'
    '\n'
    '- Goruntu isleme ile otomatik parametre cikarimi\n'
    '  => Gomululuk acisi, ramus iliskisi, derinlik seviyesi\n'
    '\n'
    '- Kural tabanli cerrahi zorluk skorlama motoru\n'
    '  => Basit / Cerrahi / Ileri Cerrahi siniflandirmasi\n'
    '\n'
    '- Web tabanli klinik karar destek arayuzu\n'
    '  => FastAPI backend + Canvas overlay gorsellesˆtirme\n'
    '\n'
    '- DICOM formati destegi (gelistirme asamasinda)\n'
    '  => Dijital radyoloji sistemleriyle dogrudan entegrasyon\n'
)

# ===================================================================
# SLIDE 12 — Proje Plani ve Is Paketleri
# ===================================================================
s = slides[11]
safe_set(s, 'PlaceHolder 3',
    'Is Paketi 1 - Veri Seti Hazirliği (Ay 1-2)\n'
    '  Panoramik rontgen goruntulerinin toplanmasi, anonimlestirilmesi\n'
    '  ve bounding box yontemiyle etiketlenmesi (%80/%10/%10 split)\n'
    '\n'
    'Is Paketi 2 - Model Gelistirme & Egitim (Ay 2-4)\n'
    '  YOLOv8/YOLO11 fine-tuning, veri zenginlestirme, post-filtre katmani\n'
    '\n'
    'Is Paketi 3 - Parametre Cikarim Algoritmasi (Ay 3-5)\n'
    '  Gomululuk acisi, ramus iliskisi ve derinlik otomatik hesaplama\n'
    '\n'
    'Is Paketi 4 - Skorlama Motoru & Arayuz (Ay 4-6)\n'
    '  Kural tabanli zorluk skoru, web arayuzu, klinik test\n'
    '\n'
    'Is Paketi 5 - DICOM Entegrasyonu & Dogrulama (Ay 6-8)\n'
    '  DICOM okuma modulu, uzman klinisyen karsilastirma calismasi\n'
)

# ===================================================================
# SLIDE 13 — Panoramik Rontgen Veri Seti & Sistem Gorunumu
# ===================================================================
s = slides[12]
safe_set(s, 'PlaceHolder 1', 'Panoramik Rontgen Veri Seti & Sistem Arayuzu')
safe_set(s, 'PlaceHolder 3',
    'Egitim veri seti:\n'
    '- Farkli gomululuk tipleri (Vertical, Mesioangular, Distoangular, Horizontal)\n'
    '- Cesitli hasta profilleri ve goruntu kalitesi kosullari\n'
    '- Bounding box etiketleme + gomululuk tipi sinifi\n'
    '- %80 egitim / %10 dogrulama / %10 test split\n'
)
safe_set(s, 'PlaceHolder 4',
    'Sistem arayuzu:\n'
    '- Surukle-birak goruntu yukleme\n'
    '- Anlik YOLO tespiti + Canvas overlay\n'
    '- Her dis icin otomatik parametre karti\n'
    '- Zorluk skoru & klinisyen oneri metni\n'
)

# ===================================================================
# SLIDE 14 — PoC Uygulamalarimiz
# ===================================================================
s = slides[13]
safe_set(s, 'PlaceHolder 1', 'Teklif Edilen Cozum : PoC Uygulamalarimiz')
safe_set(s, 'PlaceHolder 3',
    'Mevcut PoC Durumu:\n'
    '- YOLOv8/YOLO11 modeli disprojesi3 agirliklari ile egitilmis\n'
    '- Web tabanli prototip aktif: FastAPI + HTML5 Canvas overlay\n'
    '- Tespit sonrasi otomatik: aci, ramus, derinlik hesaplama\n'
    '- Anatomik post-filtre: yanlis pozitif azaltma (%90 uzeri mAP hedefi)\n'
    '- Panoramik rontgen goruntülerinde 4+ yirmilik dis tespiti\n'
)
safe_set(s, 'PlaceHolder 4',
    'Sonraki Adimlar:\n'
    '- Veri seti buyutme (daha fazla etiketli vaka)\n'
    '- DICOM formati destegi (pydicom entegrasyonu)\n'
    '- Klinik dogrulama calismasi (uzman klinisyen karsilastirmasi)\n'
    '- Mobil/tablet uyumlu arayuz gelistirme\n'
    '- Potansiyel TITCK Sinif I yazilim tibbi cihaz basvurusu\n'
)

# ===================================================================
# SLIDE 15 — Faydalar & Katma Deger
# ===================================================================
s = slides[14]
safe_set(s, 'PlaceHolder 3',
    'Klinisyen karar verme surecini standardize etme ve hizlandirma\n'
    'Onleyici degerlendirme => Sinir hasari, alveolar kirik gibi komplikasyonlarin minimize edilmesi\n'
    'Yeni mezun dis hekimleri icin uzman seviyesi on analiz destegi\n'
    'Olceklenebilirlik - ozel kliniklerden universite dis hastanelerine kadar uygulanabilir\n'
    '"Erken dogru tespit => dogru cerrahi plan => komplikasyon azalmasi, hasta memnuniyeti artisi."\n'
    '"Risk azaltimi => Sinir hasari, uzun operasyon suresi gibi sonuclarin onlenmesi."\n'
    '"Yapay zeka + kural motoru + gorsellestirme => surdurulebilir klinik karar destek altyapisi."\n'
)

# ===================================================================
# SLIDE 16 — Yenilikci Yon, Rakipler ve Rekabet Stratejisi
# ===================================================================
s = slides[15]
safe_set(s, 'PlaceHolder 3',
    'Rakip kuruluslarin farki ve rekabet stratejisi baslica birkac madde ile aciklanabilir;\n'
    'Panoramik rontgeni bir insanin yorumlamasi yerine gelistirilecek yazilim ile;\n'
    'Goruntu isleme algoritmalari ile gercek gomululuk acisi analizini yapabilmesi,\n'
    'Tespit edilen disleri numarali kutu etiketi ile pinlemesi,\n'
    'Her dis icin gomululuk tipi, ramus iliskisi ve derinlik parametrelerini\n'
    'otomatik olarak cikarabilmesi,\n'
    'Hasta profili (yas, cinsiyet, agiz acikligi) ile birlesik cerrahi zorluk skoru uretmesi,\n'
    'Zenginlestirilmis klinik raporlama yapabilmesi,\n'
    'Derin ogrenme ile rontgen goruntülerinde tespit dogrulugˆunun artirilmasi\n'
    've analizlerin daha saglikli, tekrarlanabilir bicimde yapilabilmesi.\n'
)

# ===================================================================
# SLIDE 17 — Ticarilestirme Potansiyeli
# ===================================================================
s = slides[16]
safe_set(s, 'PlaceHolder 3',
    'Musteri:\n'
    'Siparise Dayali ArGe Projesi\n'
    'Soz konusu yazilimin gelistirilmesi icin gereken tum maliyetler,\n'
    '.............. tarafindan karsilanacaktir.\n'
    'Bu kapsam; yazilim gelistirme emegi ve uretim bedelleri, test ortamlari ve\n'
    'altyapi maliyetleri, klinik dogrulama calismalari,\n'
    'lisans ve abonelik ucretleri ile proje kapsaminda gerekli gorulen\n'
    'diger giderleri icermektedir.\n'
    'Ek ArGe gereksinimi halinde ....... benzer isbirligi sozlesmesi ile siparis verecektir.\n'
    '\n'
    'Hedef Pazar:\n'
    '- Turkiye genelinde 25.000+ aktif dis hekimi\n'
    '- Ozel klinikler ve yeni mezun dis hekimleri (oncelikli segment)\n'
    '- Dis hekimligi fakulteleri (egitim kullanimi)\n'
    '- Global dis hekimligi yazilim pazari (lisanslama potansiyeli)\n'
)

# ===================================================================
# SLIDE 18 — Kaynaklar
# ===================================================================
s = slides[17]
safe_set(s, 'PlaceHolder 3',
    'Pell GJ, Gregory GT. Impacted mandibular third molars: Classification and modified\n'
    'technique for removal. Dental Digest 1933;39:330-338.\n'
    '\n'
    'Winter GB. Principles of exodontia as applied to the impacted mandibular third molar.\n'
    'American Medical Book Co; 1926.\n'
    '\n'
    'Johansson J et al. Deep learning for automated detection of impacted wisdom teeth.\n'
    'Journal of Dental Research, 2021.\n'
    '\n'
    'Ultralytics YOLOv8 Documentation\n'
    'https://docs.ultralytics.com\n'
    '\n'
    'Roboflow - Dental X-ray Detection Datasets\n'
    'https://roboflow.com\n'
    '\n'
    'pydicom - DICOM Processing in Python\n'
    'https://pydicom.github.io\n'
)

# ===================================================================
# SLIDE 19 — Tesekkurler (aynen birak)
# ===================================================================
# Degisiklik yok.

# Kaydet
prs.save(DST)
print(f'\nDosya kaydedildi: {DST}')
print('Tamamlandi!')
